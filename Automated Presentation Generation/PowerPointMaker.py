import math

from pptx import Presentation
from PIL import Image
import os
import PictureExtractorV2 as pe
import glob
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Pt

par1 = Presentation()

def makePowerPoint(textList,pageStart,firstSlideContent,fileName = "file",pictures = pe.image()):
    if len(firstSlideContent)>1:
        FirstSlide(firstSlideContent[0] , firstSlideContent[1])
    elif len(firstSlideContent) ==1:
        FirstSlide(firstSlideContent[0])
    else:
        FirstSlide()
    index = 0
    textIndex= 0
    for i in range(len(textList)):
        for j in range(index, len(pictures.imagePage)):
            if not(pictures.imagePage[j] == pageStart+i):
                index = j
                break
            makePhotoWithTexts(str(pictures.imageName[j])+'.png')
        for k in range(textIndex,len(textList)):
            if textList[k][2] == i and len(textList[k][1]) > 1:
                SlideTextsonly(textList[k][0], textList[k][1])
            else:
                textIndex = k
                break
        #if len(textList[i][1]) == 0 or textList[i][1] == ' ':
            #continue

        #listTest = spliter(textList[i])

       # length = math.celi(len(text) / 18)


    #file name
    if len(fileName) == 0:
        fileName = "file"
    par1.save(fileName+".pptx")

def getMax(num1,num2):
    if num1 > num2:
        return num1
    return num2

def FirstSlide(Text = ' ',Subtitle= ' '):
    slide1_reg = par1.slide_layouts[0]
    slide1 = par1.slides.add_slide(slide1_reg)
    title1 = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title1.text = Text
    subtitle.text = Subtitle

def makePhoto(PhotoLink):
    slide_reg = par1.slide_layouts[8]
    slide = par1.slides.add_slide(slide_reg)
    Image = slide.placeholders[1].insert_picture(PhotoLink)

def makePhotoWithTexts(PhotoLink,Text=' ',Subtext= ' '):
    slide_reg = par1.slide_layouts[8]
    slide = par1.slides.add_slide(slide_reg)
    Image = slide.placeholders[1].insert_picture(PhotoLink)
    slide.shapes.title.text=Text
    slide.placeholders[2].text=Subtext

def SlideTextsonly(Text , Subtext):
    slide_reg = par1.slide_layouts[1]
    slide = par1.slides.add_slide(slide_reg)
    title1 = slide.shapes.title
    title1.text = Text
    subtextPlaceholder = slide.placeholders[1]
   # subtextPlaceholder.font = Pt(28)
    subtextPlaceholder.text = Subtext
    #subtextPlaceholder.text.text_frame.paragraphs[0].font.size = Pt(15)



# To Load Multiple Images


def countWords(text):
    text = text.strip()
    list = text.split()
    return len(list)

def spliter(text):
    returnList = []
    count=18
    length = math.ceil(countWords(text)/18)
    print(countWords(text))
    for i in range (length):
        if countWords(text) > 18 and len(text) > 287:
            list = text.split(" ", count)
            temp = ''

            for j in range(18):
                temp += list[j]+' '
            returnList.append(temp.strip())
            count += 18
            returnList.append(list[18])
        else:
            returnList.append(text)
    return returnList
